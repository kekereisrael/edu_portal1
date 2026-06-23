"""
Exam Generation Service — Phase 6E
Intelligent auto-generation of MCQ and theory questions from uploaded materials.

Pipeline:
  1. MaterialAnalyzer  — extract text from PDF/DOCX/TXT, identify topics & keywords
  2. QuestionGenerator — call OpenAI or Anthropic to generate structured questions
  3. ExamGenerationService — orchestrate the pipeline, save results, track credits

Supported AI providers (configured via settings):
  - OpenAI  (OPENAI_API_KEY)
  - Anthropic (ANTHROPIC_API_KEY)
  - Fallback: rule-based stub questions (always works, no API key needed)
"""

import json
import logging
import re
from dataclasses import dataclass, field
from typing import Optional

from django.conf import settings

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Data Transfer Objects
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class GeneratedQuestion:
    """Represents a single generated question."""
    question_text: str
    question_type: str = 'mcq'          # 'mcq' | 'true_false' | 'theory'
    option_a: str = ''
    option_b: str = ''
    option_c: str = ''
    option_d: str = ''
    correct_answer: str = 'A'           # 'A'|'B'|'C'|'D' for MCQ; 'True'|'False' for T/F
    explanation: str = ''
    difficulty: str = 'medium'          # 'easy' | 'medium' | 'hard'
    topic: Optional[str] = None
    marks: int = 1


@dataclass
class GenerationRequest:
    """Parameters for requesting question generation."""
    subject_name: str
    topic_name: Optional[str]
    difficulty: str = 'mixed'           # 'easy' | 'medium' | 'hard' | 'mixed'
    num_questions: int = 10
    question_type: str = 'mcq'          # 'mcq' | 'theory' | 'mixed'
    material_text: Optional[str] = None
    key_topics: list = field(default_factory=list)
    keywords: list = field(default_factory=list)
    existing_questions: list = field(default_factory=list)
    language: str = 'English'


@dataclass
class GenerationResult:
    """Result of a generation attempt."""
    success: bool
    questions: list                     # List[GeneratedQuestion]
    error_message: Optional[str] = None
    tokens_used: int = 0
    model_used: str = ''
    credits_consumed: int = 0
    topics_extracted: list = field(default_factory=list)
    keywords_extracted: list = field(default_factory=list)


@dataclass
class MaterialAnalysisResult:
    """Result of analysing a study material."""
    success: bool
    extracted_text: str = ''
    key_topics: list = field(default_factory=list)
    keywords: list = field(default_factory=list)
    suggested_difficulty: str = 'medium'
    word_count: int = 0
    error_message: Optional[str] = None


# ─────────────────────────────────────────────────────────────────────────────
# Material Analyzer
# ─────────────────────────────────────────────────────────────────────────────

class MaterialAnalyzer:
    """
    Extracts text and identifies topics/keywords from uploaded study materials.

    Supported formats:
      - .txt  — direct read
      - .pdf  — PyPDF2 (if installed) or pdfplumber
      - .docx — python-docx (if installed)
      - .pptx — python-pptx (if installed)
    Falls back gracefully if libraries are not installed.
    """

    SUPPORTED_TYPES = ['document', 'image', 'link']

    # Common stop words to exclude from keyword extraction
    STOP_WORDS = {
        'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
        'of', 'with', 'by', 'from', 'is', 'are', 'was', 'were', 'be', 'been',
        'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could',
        'should', 'may', 'might', 'shall', 'can', 'this', 'that', 'these',
        'those', 'it', 'its', 'they', 'them', 'their', 'we', 'our', 'you',
        'your', 'he', 'she', 'his', 'her', 'as', 'if', 'then', 'than', 'so',
        'also', 'not', 'no', 'nor', 'yet', 'both', 'either', 'each', 'all',
        'any', 'such', 'into', 'through', 'during', 'before', 'after', 'above',
        'below', 'between', 'out', 'up', 'down', 'about', 'which', 'who',
        'what', 'when', 'where', 'how', 'why', 'there', 'here',
    }

    def analyze(self, material) -> MaterialAnalysisResult:
        """
        Analyse a Material instance and extract text + topics.

        Args:
            material: materials.models.Material instance

        Returns:
            MaterialAnalysisResult
        """
        logger.info(f'[MaterialAnalyzer] Analysing: {material.title}')

        extracted_text = ''

        # Try to extract text from file
        if material.file:
            try:
                file_path = material.file.path
                ext = file_path.lower().rsplit('.', 1)[-1] if '.' in file_path else ''

                if ext == 'txt':
                    extracted_text = self._read_txt(file_path)
                elif ext == 'pdf':
                    extracted_text = self.extract_text_from_pdf(file_path)
                elif ext == 'docx':
                    extracted_text = self.extract_text_from_docx(file_path)
                elif ext == 'pptx':
                    extracted_text = self.extract_text_from_pptx(file_path)
                else:
                    logger.info(f'[MaterialAnalyzer] Unsupported file type: {ext}')
            except Exception as e:
                logger.warning(f'[MaterialAnalyzer] File extraction failed: {e}')

        # Fall back to title + description if no file text
        if not extracted_text:
            parts = [material.title]
            if material.description:
                parts.append(material.description)
            if material.topic:
                parts.append(material.topic.name)
            extracted_text = ' '.join(parts)

        if not extracted_text.strip():
            return MaterialAnalysisResult(
                success=False,
                error_message='No text content could be extracted from this material.',
            )

        word_count = len(extracted_text.split())
        key_topics = self._extract_topics(extracted_text, material)
        keywords = self._extract_keywords(extracted_text)
        difficulty = self._suggest_difficulty(extracted_text, word_count)

        return MaterialAnalysisResult(
            success=True,
            extracted_text=extracted_text[:8000],  # Cap at 8K chars for AI prompt
            key_topics=key_topics,
            keywords=keywords[:20],
            suggested_difficulty=difficulty,
            word_count=word_count,
        )

    # ── Text extraction ───────────────────────────────────────────────────────

    def _read_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2 or pdfplumber."""
        # Try pdfplumber first (better quality)
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                return '\n'.join(
                    page.extract_text() or '' for page in pdf.pages
                )
        except ImportError:
            pass

        # Fall back to PyPDF2
        try:
            import PyPDF2
            text = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text.append(page.extract_text() or '')
            return '\n'.join(text)
        except ImportError:
            logger.warning('[MaterialAnalyzer] Neither pdfplumber nor PyPDF2 is installed.')
            return ''

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document
            doc = Document(file_path)
            return '\n'.join(para.text for para in doc.paragraphs if para.text.strip())
        except ImportError:
            logger.warning('[MaterialAnalyzer] python-docx is not installed.')
            return ''

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text.append(shape.text)
            return '\n'.join(text)
        except ImportError:
            logger.warning('[MaterialAnalyzer] python-pptx is not installed.')
            return ''

    # ── NLP helpers ───────────────────────────────────────────────────────────

    def _extract_topics(self, text: str, material) -> list:
        """Extract key topics from text. Uses topic FK if available."""
        topics = []

        # Use the material's linked topic first
        if material.topic:
            topics.append(material.topic.name)

        # Extract capitalized multi-word phrases (likely topic names)
        # Pattern: 2-4 consecutive capitalized words
        pattern = r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,3})\b'
        matches = re.findall(pattern, text)
        # Deduplicate and filter short matches
        seen = set(topics)
        for m in matches:
            if m not in seen and len(m) > 5:
                topics.append(m)
                seen.add(m)
                if len(topics) >= 10:
                    break

        # Also extract section headers (lines ending with colon or all-caps lines)
        for line in text.split('\n'):
            line = line.strip()
            if line.endswith(':') and 3 < len(line) < 60:
                topic = line.rstrip(':').strip()
                if topic not in seen:
                    topics.append(topic)
                    seen.add(topic)
            if len(topics) >= 10:
                break

        return topics[:10]

    def _extract_keywords(self, text: str) -> list:
        """Extract significant keywords using word frequency analysis."""
        # Tokenize and clean
        words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
        # Remove stop words
        words = [w for w in words if w not in self.STOP_WORDS]

        # Count frequency
        freq = {}
        for word in words:
            freq[word] = freq.get(word, 0) + 1

        # Sort by frequency, return top 20
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [word for word, _ in sorted_words[:20]]

    def _suggest_difficulty(self, text: str, word_count: int) -> str:
        """Suggest difficulty based on text complexity."""
        # Simple heuristic: average word length + sentence length
        words = text.split()
        if not words:
            return 'medium'
        avg_word_len = sum(len(w) for w in words) / len(words)
        sentences = re.split(r'[.!?]+', text)
        avg_sent_len = len(words) / max(len(sentences), 1)

        if avg_word_len > 7 or avg_sent_len > 25:
            return 'hard'
        elif avg_word_len < 5 and avg_sent_len < 12:
            return 'easy'
        return 'medium'


# ─────────────────────────────────────────────────────────────────────────────
# Question Generator
# ─────────────────────────────────────────────────────────────────────────────

class QuestionGenerator:
    """
    Generates exam questions using AI (OpenAI or Anthropic).
    Falls back to rule-based stub generation if no API key is configured.
    """

    # JSON schema the AI must return
    MCQ_SCHEMA = """
Return a JSON array of objects. Each object must have exactly these fields:
{
  "question_text": "...",
  "question_type": "mcq",
  "option_a": "...",
  "option_b": "...",
  "option_c": "...",
  "option_d": "...",
  "correct_answer": "A",
  "explanation": "...",
  "difficulty": "easy|medium|hard",
  "topic": "..."
}
"""

    THEORY_SCHEMA = """
Return a JSON array of objects. Each object must have exactly these fields:
{
  "question_text": "...",
  "question_type": "theory",
  "option_a": "",
  "option_b": "",
  "option_c": "",
  "option_d": "",
  "correct_answer": "",
  "explanation": "Model answer or marking guide",
  "difficulty": "easy|medium|hard",
  "topic": "..."
}
"""

    def generate(self, request: GenerationRequest) -> GenerationResult:
        """
        Generate questions. Tries OpenAI → Anthropic → stub fallback.
        """
        logger.info(
            f'[QuestionGenerator] Generating {request.num_questions} '
            f'{request.question_type} questions for {request.subject_name}'
        )

        # Try OpenAI
        if getattr(settings, 'OPENAI_API_KEY', ''):
            result = self._generate_openai(request)
            if result.success:
                return result
            logger.warning(f'[QuestionGenerator] OpenAI failed: {result.error_message}')

        # Try Anthropic
        if getattr(settings, 'ANTHROPIC_API_KEY', ''):
            result = self._generate_anthropic(request)
            if result.success:
                return result
            logger.warning(f'[QuestionGenerator] Anthropic failed: {result.error_message}')

        # Fallback: stub questions
        logger.info('[QuestionGenerator] Using stub fallback (no AI key configured)')
        return self._generate_stub(request)

    # ── OpenAI ────────────────────────────────────────────────────────────────

    def _generate_openai(self, request: GenerationRequest) -> GenerationResult:
        try:
            import openai
            client = openai.OpenAI(api_key=settings.OPENAI_API_KEY)

            prompt = self._build_prompt(request)
            response = client.chat.completions.create(
                model='gpt-4o-mini',
                messages=[
                    {
                        'role': 'system',
                        'content': (
                            'You are an expert educational assessment designer. '
                            'Generate high-quality exam questions in valid JSON only. '
                            'Return ONLY the JSON array, no markdown, no explanation.'
                        ),
                    },
                    {'role': 'user', 'content': prompt},
                ],
                temperature=0.7,
                max_tokens=4000,
            )

            raw = response.choices[0].message.content.strip()
            questions = self._parse_ai_response(raw)
            tokens = response.usage.total_tokens if response.usage else 0

            return GenerationResult(
                success=True,
                questions=questions,
                model_used='gpt-4o-mini',
                tokens_used=tokens,
                credits_consumed=max(1, tokens // 1000),
            )

        except ImportError:
            return GenerationResult(
                success=False, questions=[],
                error_message='openai package not installed. Run: pip install openai',
            )
        except Exception as e:
            return GenerationResult(
                success=False, questions=[],
                error_message=f'OpenAI error: {e}',
            )

    # ── Anthropic ─────────────────────────────────────────────────────────────

    def _generate_anthropic(self, request: GenerationRequest) -> GenerationResult:
        try:
            import anthropic
            client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)

            prompt = self._build_prompt(request)
            message = client.messages.create(
                model='claude-3-haiku-20240307',
                max_tokens=4000,
                messages=[{'role': 'user', 'content': prompt}],
                system=(
                    'You are an expert educational assessment designer. '
                    'Generate high-quality exam questions in valid JSON only. '
                    'Return ONLY the JSON array, no markdown, no explanation.'
                ),
            )

            raw = message.content[0].text.strip()
            questions = self._parse_ai_response(raw)
            tokens = (message.usage.input_tokens + message.usage.output_tokens
                      if message.usage else 0)

            return GenerationResult(
                success=True,
                questions=questions,
                model_used='claude-3-haiku',
                tokens_used=tokens,
                credits_consumed=max(1, tokens // 1000),
            )

        except ImportError:
            return GenerationResult(
                success=False, questions=[],
                error_message='anthropic package not installed. Run: pip install anthropic',
            )
        except Exception as e:
            return GenerationResult(
                success=False, questions=[],
                error_message=f'Anthropic error: {e}',
            )

    # ── Stub fallback ─────────────────────────────────────────────────────────

    def _generate_stub(self, request: GenerationRequest) -> GenerationResult:
        """
        Generate placeholder questions when no AI key is available.
        Useful for development/testing.
        """
        questions = []
        topic = request.topic_name or request.subject_name
        keywords = request.keywords[:5] if request.keywords else [topic]

        difficulties = ['easy', 'medium', 'hard']

        for i in range(request.num_questions):
            diff = difficulties[i % 3] if request.difficulty == 'mixed' else request.difficulty
            kw = keywords[i % len(keywords)]

            if request.question_type in ('mcq', 'mixed') and (
                request.question_type == 'mcq' or i % 3 != 2
            ):
                q = GeneratedQuestion(
                    question_text=f'Which of the following best describes {kw} in the context of {topic}?',
                    question_type='mcq',
                    option_a=f'The primary definition of {kw}',
                    option_b=f'An unrelated concept to {kw}',
                    option_c=f'A common misconception about {kw}',
                    option_d=f'An advanced application of {kw}',
                    correct_answer='A',
                    explanation=f'Option A correctly defines {kw} as covered in the material.',
                    difficulty=diff,
                    topic=topic,
                    marks=1 if diff == 'easy' else (2 if diff == 'medium' else 3),
                )
            else:
                q = GeneratedQuestion(
                    question_text=f'Explain the concept of {kw} and its significance in {topic}.',
                    question_type='theory',
                    explanation=f'A complete answer should define {kw}, explain its role in {topic}, and provide at least one example.',
                    difficulty=diff,
                    topic=topic,
                    marks=5 if diff == 'easy' else (10 if diff == 'medium' else 15),
                )
            questions.append(q)

        return GenerationResult(
            success=True,
            questions=questions,
            model_used='stub',
            tokens_used=0,
            credits_consumed=0,
        )

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _build_prompt(self, request: GenerationRequest) -> str:
        """Build the AI prompt from the generation request."""
        schema = self.MCQ_SCHEMA if request.question_type != 'theory' else self.THEORY_SCHEMA

        context_parts = [f'Subject: {request.subject_name}']
        if request.topic_name:
            context_parts.append(f'Topic: {request.topic_name}')
        if request.key_topics:
            context_parts.append(f'Key topics: {", ".join(request.key_topics[:5])}')
        if request.keywords:
            context_parts.append(f'Keywords: {", ".join(request.keywords[:10])}')

        difficulty_instruction = (
            f'Mix difficulties (roughly 30% easy, 50% medium, 20% hard)'
            if request.difficulty == 'mixed'
            else f'All questions should be {request.difficulty} difficulty'
        )

        material_section = ''
        if request.material_text:
            material_section = f'\n\nStudy material excerpt:\n"""\n{request.material_text[:3000]}\n"""'

        type_instruction = {
            'mcq': f'Generate {request.num_questions} multiple-choice questions (4 options each).',
            'theory': f'Generate {request.num_questions} theory/essay questions.',
            'mixed': (
                f'Generate {request.num_questions} questions: '
                f'{request.num_questions * 2 // 3} MCQ and '
                f'{request.num_questions // 3} theory questions.'
            ),
        }.get(request.question_type, f'Generate {request.num_questions} questions.')

        return f"""{type_instruction}

Context:
{chr(10).join(context_parts)}

Difficulty: {difficulty_instruction}
Language: {request.language}
{material_section}

{schema}

Important:
- Questions must be factually accurate and educationally appropriate
- MCQ options must be plausible (no obviously wrong answers)
- Explanations must be clear and educational
- Return ONLY the JSON array, nothing else
"""

    def _parse_ai_response(self, response_text: str) -> list:
        """Parse AI JSON response into GeneratedQuestion objects."""
        # Strip markdown code blocks if present
        text = response_text.strip()
        if text.startswith('```'):
            text = re.sub(r'^```(?:json)?\s*', '', text)
            text = re.sub(r'\s*```$', '', text)

        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            # Try to extract JSON array from response
            match = re.search(r'\[.*\]', text, re.DOTALL)
            if match:
                data = json.loads(match.group())
            else:
                raise ValueError(f'Could not parse JSON from AI response: {text[:200]}')

        questions = []
        for item in data:
            if not isinstance(item, dict):
                continue
            q = GeneratedQuestion(
                question_text=item.get('question_text', '').strip(),
                question_type=item.get('question_type', 'mcq'),
                option_a=item.get('option_a', ''),
                option_b=item.get('option_b', ''),
                option_c=item.get('option_c', ''),
                option_d=item.get('option_d', ''),
                correct_answer=item.get('correct_answer', 'A').upper(),
                explanation=item.get('explanation', ''),
                difficulty=item.get('difficulty', 'medium'),
                topic=item.get('topic'),
                marks=int(item.get('marks', 1)),
            )
            if self._validate_question(q):
                questions.append(q)

        return questions

    def _validate_question(self, question: GeneratedQuestion) -> bool:
        """Validate a generated question for completeness."""
        if not question.question_text or len(question.question_text) < 10:
            return False
        if question.question_type == 'mcq':
            if question.correct_answer not in ('A', 'B', 'C', 'D'):
                return False
            if not all([question.option_a, question.option_b,
                        question.option_c, question.option_d]):
                return False
        return True


# ─────────────────────────────────────────────────────────────────────────────
# High-Level Service Facade
# ─────────────────────────────────────────────────────────────────────────────

class ExamGenerationService:
    """
    Orchestrates material analysis and question generation.
    Use this class in views, signals, and Celery tasks.
    """

    def __init__(self):
        self.analyzer = MaterialAnalyzer()
        self.generator = QuestionGenerator()

    def generate_from_material(
        self,
        material,
        num_questions: int = 10,
        difficulty: str = 'mixed',
        question_type: str = 'mcq',
        school=None,
        user=None,
    ) -> GenerationResult:
        """
        Generate exam questions from an uploaded study material.

        Args:
            material: materials.models.Material instance
            num_questions: How many questions to generate
            difficulty: 'easy' | 'medium' | 'hard' | 'mixed'
            question_type: 'mcq' | 'theory' | 'mixed'
            school: School instance (for credit tracking)
            user: User instance (for credit tracking)

        Returns:
            GenerationResult
        """
        # Step 1: Analyse material
        analysis = self.analyzer.analyze(material)
        if not analysis.success:
            return GenerationResult(
                success=False,
                questions=[],
                error_message=f'Material analysis failed: {analysis.error_message}',
            )

        logger.info(
            f'[ExamGenService] Material analysed: {analysis.word_count} words, '
            f'{len(analysis.key_topics)} topics, {len(analysis.keywords)} keywords'
        )

        # Step 2: Build generation request
        request = GenerationRequest(
            subject_name=material.subject.name if material.subject else 'General',
            topic_name=material.topic.name if material.topic else None,
            difficulty=difficulty or analysis.suggested_difficulty,
            num_questions=num_questions,
            question_type=question_type,
            material_text=analysis.extracted_text,
            key_topics=analysis.key_topics,
            keywords=analysis.keywords,
        )

        # Step 3: Generate questions
        result = self.generator.generate(request)
        result.topics_extracted = analysis.key_topics
        result.keywords_extracted = analysis.keywords

        # Step 4: Track AI usage
        if result.success and result.credits_consumed > 0 and school and user:
            self._track_usage(school, user, result)

        return result

    def generate_from_topic(
        self,
        subject_name: str,
        topic_name: str,
        num_questions: int = 10,
        difficulty: str = 'medium',
        question_type: str = 'mcq',
    ) -> GenerationResult:
        """
        Generate questions for a specific topic without a material file.
        """
        request = GenerationRequest(
            subject_name=subject_name,
            topic_name=topic_name,
            difficulty=difficulty,
            num_questions=num_questions,
            question_type=question_type,
        )
        return self.generator.generate(request)

    def save_generated_questions(
        self,
        exam,
        result: GenerationResult,
        question_bank=None,
    ) -> int:
        """
        Save generated questions to an Exam or QuestionBank.

        Args:
            exam: exams.models.Exam instance (or None if saving to bank only)
            result: GenerationResult with questions
            question_bank: exams.models.QuestionBank instance (optional)

        Returns:
            Number of questions saved
        """
        if not result.success or not result.questions:
            return 0

        saved = 0

        if exam:
            from exams.models import Question
            last_order = exam.questions.count()
            for i, q in enumerate(result.questions):
                if not self.generator._validate_question(q):
                    continue
                Question.objects.create(
                    exam=exam,
                    question_text=q.question_text,
                    question_type=q.question_type,
                    option_a=q.option_a,
                    option_b=q.option_b,
                    option_c=q.option_c,
                    option_d=q.option_d,
                    correct_answer=q.correct_answer,
                    explanation=q.explanation,
                    difficulty=q.difficulty,
                    marks=q.marks,
                    order=last_order + i + 1,
                )
                saved += 1
            if saved > 0:
                try:
                    exam.recalculate_total_marks()
                except Exception:
                    pass

        if question_bank:
            from exams.models import BankQuestion
            for q in result.questions:
                if not self.generator._validate_question(q):
                    continue
                try:
                    BankQuestion.objects.create(
                        bank=question_bank,
                        question_text=q.question_text,
                        question_type=q.question_type,
                        option_a=q.option_a,
                        option_b=q.option_b,
                        option_c=q.option_c,
                        option_d=q.option_d,
                        correct_answer=q.correct_answer,
                        explanation=q.explanation,
                        difficulty=q.difficulty,
                        marks=q.marks,
                        topic=self._resolve_topic(q.topic, question_bank),
                    )
                    if not exam:
                        saved += 1
                except Exception as e:
                    logger.warning(f'[ExamGenService] Failed to save bank question: {e}')

        logger.info(f'[ExamGenService] Saved {saved} questions')
        return saved

    def _resolve_topic(self, topic_name: str, question_bank):
        """Try to find a matching Topic FK for the bank question."""
        if not topic_name or not question_bank:
            return None
        try:
            from subjects.models import Topic
            return Topic.objects.filter(
                name__iexact=topic_name,
                subject=question_bank.subject,
            ).first()
        except Exception:
            return None

    def _track_usage(self, school, user, result: GenerationResult):
        """Track AI credit usage in analytics.models.AIUsageRecord."""
        try:
            from analytics.models import AIUsageRecord
            AIUsageRecord.objects.create(
                school=school,
                user=user,
                usage_type=AIUsageRecord.UsageType.QUESTION_GENERATION,
                credits_consumed=result.credits_consumed,
                input_tokens=result.tokens_used // 2,
                output_tokens=result.tokens_used // 2,
                model_used=result.model_used or 'unknown',
            )
        except Exception as e:
            logger.warning(f'[ExamGenService] Failed to track AI usage: {e}')


# ─────────────────────────────────────────────────────────────────────────────
# Module-level singleton
# ─────────────────────────────────────────────────────────────────────────────

exam_generation_service = ExamGenerationService()
